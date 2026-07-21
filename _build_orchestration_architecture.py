"""
Generate a standalone, single-slide architecture diagram:

    "Order Creation & Orchestration - Architecture"

The diagram is a layered reference view (left tier labels + cross-cutting rail)
covering, end to end:
  - Upstream / customer interaction: the 4 PO intake methods + ingestion
  - Experience: CSR workspace & console
  - Orchestration: resumable pipeline/state machine + auto-approval and CSR paths
  - Order creation & downstream: order execution -> ERP / OMS / WMS / carrier / notify
  - Data foundation: governed master data (the data process behind every decision)
  - Cross-cutting: exception governance, audit, security, observability

Output: demo/Order-Creation-Orchestration-Architecture.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ── Palette (EY-inspired: near-black + signature yellow) ──────────────────────
BG     = RGBColor(0x14, 0x14, 0x1E)
BG2    = RGBColor(0x0E, 0x0E, 0x16)
PANEL  = RGBColor(0x24, 0x24, 0x33)
PANEL2 = RGBColor(0x2E, 0x2E, 0x40)
YELLOW = RGBColor(0xFF, 0xE6, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0xB8, 0xB8, 0xC6)
DGREY  = RGBColor(0x8A, 0x8A, 0x99)
GREEN  = RGBColor(0x35, 0xC7, 0x59)
RED    = RGBColor(0xFF, 0x5A, 0x5A)
AMBER  = RGBColor(0xFF, 0xB0, 0x20)
BLUE   = RGBColor(0x4A, 0x9E, 0xFF)
TEAL   = RGBColor(0x33, 0xC9, 0xC9)
PURPLE = RGBColor(0xA9, 0x7B, 0xFF)
LINE   = RGBColor(0x3C, 0x3C, 0x4E)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def _fill(shape, color, line_color=None, line_w=None):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color; shape.line.width = Pt(line_w or 1)
    shape.shadow.inherit = False


def box(s, x, y, w, h, fill=PANEL, line_color=None, line_w=None, radius=True, adj=0.06):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(shp, fill, line_color, line_w)
    if radius:
        try: shp.adjustments[0] = adj
        except Exception: pass
    shp.text_frame.word_wrap = True
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=1):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0); p.line_spacing = line_spacing
        if isinstance(para, tuple): para = [para]
        for (t, sz, col, bold, ital) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold; r.font.italic = ital; r.font.name = FONT
    return tb


def fill_text(shape, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05); tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align; p.space_after = Pt(1); p.space_before = Pt(0)
        if isinstance(line, tuple): line = [line]
        for (t, sz, col, bold, ital) in line:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold; r.font.italic = ital; r.font.name = FONT


def tile(s, x, y, w, h, lines, fill=PANEL2, line_color=LINE, line_w=0.75, fs=9, tcol=WHITE):
    b = box(s, x, y, w, h, fill=fill, line_color=line_color, line_w=line_w)
    fill_text(b, [[(ln, fs, tcol, True, False)] for ln in lines])
    return b


def _lprops(conn, color, w, dash=None, tail=False, head=False):
    ln = conn.line; ln.color.rgb = color; ln.width = Pt(w)
    el = ln._get_or_add_ln()
    if dash: el.append(el.makeelement(qn('a:prstDash'), {'val': dash}))
    if head: el.append(el.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if tail: el.append(el.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))


def conn(s, x1, y1, x2, y2, color=YELLOW, w=1.5, dash=None, tail=True, head=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.shadow.inherit = False; _lprops(c, color, w, dash, tail, head)
    return c


def chevron(s, x, y, w, h=0.24, color=YELLOW):
    a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(a, color); return a


def down(s, x, y, w=0.3, h=0.24, color=YELLOW):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(a, color); return a


def label(s, x, y, w, text, color=YELLOW, size=8.5, align=PP_ALIGN.CENTER):
    txt(s, x, y, w, 0.2, [[(text, size, color, True, False)]], align=align)


# ── Slide ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bgr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
_fill(bgr, BG); bgr.text_frame.paragraphs[0].text = ""

# title bar
box(s, 0, 0, 13.333, 0.12, fill=YELLOW, radius=False)
box(s, 0, 0.12, 0.16, 0.86, fill=YELLOW, radius=False)
txt(s, 0.45, 0.14, 9.5, 0.4, [[("SOLUTION ARCHITECTURE", 11, YELLOW, True, False)]])
txt(s, 0.45, 0.38, 12.5, 0.55,
    [[("Order Creation & Orchestration \u2014 Architecture ", 22, WHITE, True, False),
      ("with Human-in-the-Loop Governance", 22, YELLOW, True, False)]])

# geometry
LX, LW = 0.4, 1.5           # tier-label column
MX, MW = 2.0, 9.0           # main content (2.0 .. 11.0)
XR, XRW = 11.18, 1.82       # cross-cutting rail (11.18 .. 13.0)
MCX = MX + MW / 2

# band geometry: (y, h, label, color)
L1 = (1.15, 1.05, "UPSTREAM\nINTAKE", BLUE)
L2 = (2.28, 0.50, "EXPERIENCE", TEAL)
L3 = (2.86, 2.05, "ORCHESTRATION", YELLOW)
L4 = (4.99, 1.00, "ORDER CREATION\nDOWNSTREAM", GREEN)
L5 = (6.07, 0.85, "DATA\nFOUNDATION", PURPLE)


def tier(band):
    y, h, name, color = band
    b = box(s, LX, y, LW, h, fill=PANEL)
    box(s, LX, y, 0.07, h, fill=color, radius=False)
    fill_text(b, [[(w, 9, WHITE, True, False)] for w in name.split("\n")])


for b in (L1, L2, L3, L4, L5):
    tier(b)

# ── L1 UPSTREAM: 4 PO methods + ingestion ────────────────────────────────────
y, h = L1[0], L1[1]
txt(s, MX + 0.05, y + 0.05, 4.6, 0.22, [[("4 PO INTAKE METHODS  \u00b7  customer interaction", 9, YELLOW, True, False)]])
po_methods = [("Email PO", "\u2709"), ("PDF PO", "\U0001F4C4"), ("Excel PO", "\U0001F4CA"), ("Scanned PO", "\U0001F5B6")]
cw = 1.06; g = 0.1; x0 = MX + 0.05
for i, (nm, ic) in enumerate(po_methods):
    tile(s, x0 + i * (cw + g), y + 0.34, cw, 0.58, [ic + "  " + nm], fill=PANEL2, line_color=BLUE, line_w=0.9, fs=9, tcol=WHITE)
gx = x0 + 4 * (cw + g) - 0.02
chevron(s, gx, y + 0.5, 0.42, 0.26, color=BLUE)
ing = box(s, gx + 0.5, y + 0.30, 3.55, 0.62, fill=PANEL2, line_color=YELLOW, line_w=1)
fill_text(ing, [[("INGESTION & EXTRACTION", 9.5, WHITE, True, False)],
                [("OCR \u00b7 parse \u00b7 LLM extract \u00b7 normalize \u00b7 validate schema", 7.8, GREY, False, False)]])

# ── L2 EXPERIENCE: CSR workspace ─────────────────────────────────────────────
y, h = L2[0], L2[1]
eb = box(s, MX, y, MW, h, fill=PANEL2, line_color=TEAL, line_w=1)
fill_text(eb, [[("EXPERIENCE LAYER  \u2014  CSR Workspace & Console", 10.5, WHITE, True, False),
                ("     PO intake \u00b7 live agent panels \u00b7 one-click decisions \u00b7 audit viewer", 8.8, GREY, False, True)]])

# ── L3 ORCHESTRATION ─────────────────────────────────────────────────────────
y, h = L3[0], L3[1]
box(s, MX, y, MW, h, fill=BG2, line_color=LINE, line_w=0.75)
txt(s, MX + 0.12, y + 0.06, MW - 0.24, 0.25,
    [[("ORCHESTRATION ENGINE  ", 10, YELLOW, True, False),
      ("\u2014  resumable pipeline / state machine  (pause on exception \u00b7 resume on decision)", 8.8, GREY, False, True)]])
# pipeline tiles
stages = ["Intake", "Customer\nValidation", "Product\nMatching", "Pricing &\nPromo", "Credit", "Inventory", "Shipments", "Approvals"]
n = len(stages); pg = 0.1; pw = (MW - 0.24 - pg * (n - 1)) / n; px0 = MX + 0.12; pty = y + 0.42
for i, st in enumerate(stages):
    b = box(s, px0 + i * (pw + pg), pty, pw, 0.62, fill=PANEL2, line_color=YELLOW, line_w=0.8)
    fill_text(b, [[(str(i + 1), 8, YELLOW, True, False)]] + [[(ln, 8, WHITE, True, False)] for ln in st.split("\n")])
    if i < n - 1:
        chevron(s, px0 + i * (pw + pg) + pw - 0.03, pty + 0.19, pg + 0.05, 0.22, color=YELLOW)
# two processing paths
paths_y = y + 1.28; ph = 0.68
half = (MW - 0.24 - 0.2) / 2
# auto approval (green)
ab = box(s, MX + 0.12, paths_y, half, ph, fill=PANEL, line_color=GREEN, line_w=1.1)
box(s, MX + 0.12, paths_y, 0.08, ph, fill=GREEN, radius=False)
txt(s, MX + 0.3, paths_y + 0.06, half - 0.3, 0.24, [[("\u2713  AUTO-APPROVAL  \u00b7  straight-through", 9.5, GREEN, True, False)]])
txt(s, MX + 0.3, paths_y + 0.32, half - 0.3, 0.32,
    [[("Clean PO \u2014 every gate passes automatically. Zero human touch \u2192 order creation.", 8.3, GREY, False, False)]], line_spacing=1.0)
# csr (amber)
cbx = MX + 0.12 + half + 0.2
cb = box(s, cbx, paths_y, half, ph, fill=PANEL, line_color=AMBER, line_w=1.1)
box(s, cbx, paths_y, 0.08, ph, fill=AMBER, radius=False)
txt(s, cbx + 0.18, paths_y + 0.06, half - 0.3, 0.24, [[("\u26A0  CSR  \u00b7  human-in-the-loop", 9.5, AMBER, True, False)]])
txt(s, cbx + 0.18, paths_y + 0.32, half - 0.3, 0.32,
    [[("Exception \u2192 CSR resolves / approves in workspace \u2192 pipeline resumes  (reject \u2192 routed).", 8.3, GREY, False, False)]], line_spacing=1.0)
# branch labels from pipeline to paths
conn(s, MX + 0.12 + half / 2, pty + 0.62, MX + 0.12 + half / 2, paths_y, color=GREEN, w=1.2, tail=True)
label(s, MX + 0.12 + half / 2 - 0.9, pty + 0.63, 1.8, "clean \u2192 auto", GREEN, 8)
conn(s, cbx + half / 2, pty + 0.62, cbx + half / 2, paths_y, color=AMBER, w=1.2, tail=True)
label(s, cbx + half / 2 - 0.9, pty + 0.63, 1.8, "exception \u2192 CSR", AMBER, 8)
# resume loop (amber, back to pipeline)
conn(s, cbx + half, paths_y + ph / 2, cbx + half + 0.18, paths_y + ph / 2, color=AMBER, w=1.1, dash='dash', tail=False)
conn(s, cbx + half + 0.18, paths_y + ph / 2, cbx + half + 0.18, pty + 0.31, color=AMBER, w=1.1, dash='dash', tail=False)
conn(s, cbx + half + 0.18, pty + 0.31, px0 + (n - 1) * (pw + pg) + pw, pty + 0.31, color=AMBER, w=1.1, dash='dash', tail=True)
label(s, cbx + half - 1.2, paths_y - 0.02, 1.7, "resume", AMBER, 7.5, align=PP_ALIGN.RIGHT)

# ── L4 ORDER CREATION + DOWNSTREAM ───────────────────────────────────────────
y, h = L4[0], L4[1]
oe = box(s, MX + 0.05, y + 0.24, 2.55, 0.62, fill=PANEL2, line_color=GREEN, line_w=1.1)
fill_text(oe, [[("ORDER EXECUTION", 9.8, WHITE, True, False)], [("create sales order", 8.3, GREY, False, False)]])
chevron(s, MX + 2.62, y + 0.42, 0.42, 0.26, color=GREEN)
ds = ["ERP / OMS", "WMS /\nFulfillment", "Carrier /\nShipping", "Invoicing /\nFinance", "Notify +\nTracking"]
dx0 = MX + 3.15; dn = len(ds); dg = 0.1; dw = (MW - 3.15 - 0.05 - dg * (dn - 1)) / dn
for i, d in enumerate(ds):
    b = box(s, dx0 + i * (dw + dg), y + 0.24, dw, 0.62, fill=PANEL2, line_color=LINE, line_w=0.8)
    fill_text(b, [[(ln, 8.3, WHITE, True, False)] for ln in d.split("\n")])
txt(s, MX + 3.15, y + 0.02, MW - 3.15, 0.2, [[("DOWNSTREAM SYSTEMS", 8.5, YELLOW, True, False)]])

# ── L5 DATA FOUNDATION ───────────────────────────────────────────────────────
y, h = L5[0], L5[1]
box(s, MX, y, MW, h, fill=BG2, line_color=PURPLE, line_w=0.9)
txt(s, MX + 0.12, y + 0.06, MW - 0.24, 0.2,
    [[("MASTER-DATA FOUNDATION (governed)", 9, YELLOW, True, False),
      ("   \u2014  the data process behind every decision", 8.3, GREY, False, True)]])
ent = ["Customer", "Buyer", "Product", "Pricing", "Credit", "Inventory", "Logistics", "Budget", "Compliance", "Governance"]
en = len(ent); eg = 0.08; ew = (MW - 0.24 - eg * (en - 1)) / en; ex0 = MX + 0.12; ety = y + 0.36
for i, e in enumerate(ent):
    b = box(s, ex0 + i * (ew + eg), ety, ew, 0.36, fill=PANEL, line_color=LINE, line_w=0.6)
    fill_text(b, [[(e, 7.6, GREY, True, False)]])

# ── vertical flow arrows down the main spine ─────────────────────────────────
for gy in (L1[0] + L1[1], L2[0] + L2[1], L3[0] + L3[1]):
    down(s, MCX - 0.15, gy - 0.02, 0.3, 0.14, color=YELLOW)
# data feeds decisions (dashed up the left gutter)
conn(s, MX - 0.06, L5[0], MX - 0.06, L3[0] + L3[1], color=PURPLE, w=1.1, dash='dash', tail=True)
txt(s, LX + 0.02, (L4[0] + L4[1] / 2) - 0.1, LW, 0.2, [[("data \u2191", 7.5, PURPLE, True, True)]], align=PP_ALIGN.RIGHT)

# ── cross-cutting rail ───────────────────────────────────────────────────────
cc_y = L1[0]; cc_h = (L4[0] + L4[1]) - L1[0]
box(s, XR, cc_y, XRW, cc_h, fill=PANEL)
box(s, XR, cc_y, XRW, 0.1, fill=YELLOW, radius=False)
txt(s, XR + 0.12, cc_y + 0.14, XRW - 0.24, 0.25, [[("CROSS-CUTTING", 9, YELLOW, True, False)]])
cross = [("Exception\nGovernance & routing", RED), ("Audit &\nTraceability", BLUE),
         ("Security &\nAccess control", TEAL), ("Observability\n& Logging", GREEN)]
inner = cc_h - 0.5; ch = (inner - 0.3) / 4
for i, (t, col) in enumerate(cross):
    yy = cc_y + 0.5 + i * (ch + 0.1)
    b = box(s, XR + 0.12, yy, XRW - 0.24, ch, fill=PANEL2, line_color=LINE, line_w=0.7)
    box(s, XR + 0.12, yy, 0.06, ch, fill=col, radius=False)
    fill_text(b, [[(ln, 8.3, WHITE, True, False)] for ln in t.split("\n")])

# ── footer legend ────────────────────────────────────────────────────────────
txt(s, 0.45, 7.06, 12.4, 0.3,
    [[("Flow: ", 8.5, DGREY, True, True),
      ("solid = order flow   \u00b7   ", 8.5, YELLOW, False, True),
      ("dashed amber = CSR resume loop   \u00b7   ", 8.5, AMBER, False, True),
      ("dashed purple = master-data reads   \u00b7   ", 8.5, PURPLE, False, True),
      ("green = straight-through / auto-approval", 8.5, GREEN, False, True)]])

out = os.path.join("demo", "Order-Creation-Orchestration-Architecture.pptx")
prs.save(out)
print("Saved:", out, "slides:", len(prs.slides._sldIdLst))
