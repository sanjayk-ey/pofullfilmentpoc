"""
Generate a client-ready architecture slide deck for the
Order Assistant AI — PO-to-Fulfillment Orchestration Accelerator.

Structure:
  1. Business context  (need / problem, solution, impact)
  2. Solution architecture & design  (reference architecture, component view,
     process flowchart, sequence diagram)
  3. The AI tool we built  (agents, master data, pipeline, demos, governance,
     technology)

Output: demo/PO-Fulfillment-Orchestration-Architecture.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ── Palette (EY-inspired: near-black + signature yellow) ──────────────────────
BG       = RGBColor(0x14, 0x14, 0x1E)
BG2      = RGBColor(0x0E, 0x0E, 0x16)
PANEL    = RGBColor(0x24, 0x24, 0x33)
PANEL2   = RGBColor(0x2E, 0x2E, 0x40)
YELLOW   = RGBColor(0xFF, 0xE6, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY     = RGBColor(0xB8, 0xB8, 0xC6)
DGREY    = RGBColor(0x8A, 0x8A, 0x99)
GREEN    = RGBColor(0x35, 0xC7, 0x59)
RED      = RGBColor(0xFF, 0x5A, 0x5A)
AMBER    = RGBColor(0xFF, 0xB0, 0x20)
BLUE     = RGBColor(0x4A, 0x9E, 0xFF)
TEAL     = RGBColor(0x2FD, 0x0, 0x0) if False else RGBColor(0x33, 0xC9, 0xC9)
PURPLE   = RGBColor(0xA9, 0x7B, 0xFF)
LINE     = RGBColor(0x3C, 0x3C, 0x4E)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Core helpers ──────────────────────────────────────────────────────────────
_IDX = [0]
_SEC = [0]


def slide(bg=BG):
    _IDX[0] += 1
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    r.text_frame.paragraphs[0].text = ""
    return s


def _set_fill(shape, color, line_color=None, line_w=None):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w or 1)
    shape.shadow.inherit = False


def box(s, x, y, w, h, fill=PANEL, line_color=None, line_w=None, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    _set_fill(shp, fill, line_color, line_w)
    if radius:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    shp.text_frame.word_wrap = True
    return shp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=2, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after)
        p.space_before = Pt(0); p.line_spacing = line_spacing
        if isinstance(para, tuple):
            para = [para]
        for (t, sz, col, bold, ital) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.italic = ital; r.font.name = FONT
    return tb


def fill_text(shape, lines, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(1); p.space_before = Pt(0)
        if isinstance(line, tuple):
            line = [line]
        for (t, sz, col, bold, ital) in line:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.italic = ital; r.font.name = FONT


def accent_bar(s, x=0.0, y=0.0, w=0.16, h=7.5):
    box(s, x, y, w, h, fill=YELLOW, radius=False)


def header(s, kicker, title):
    accent_bar(s)
    txt(s, 0.6, 0.42, 11.5, 0.4, [[(kicker.upper(), 12, YELLOW, True, False)]])
    txt(s, 0.6, 0.72, 12.2, 0.9, [[(title, 29, WHITE, True, False)]])
    box(s, 0.62, 1.55, 2.2, 0.03, fill=YELLOW, radius=False)
    txt(s, 12.2, 6.95, 0.9, 0.4, [[(f"{_IDX[0]:02d}", 10, DGREY, False, False)]],
        align=PP_ALIGN.RIGHT)
    txt(s, 0.6, 6.95, 9.0, 0.4,
        [[("Order Assistant AI  ·  PO-to-Fulfillment Orchestration", 9, DGREY, False, False)]])


def chip(s, x, y, w, h, label, color, txtcolor=None, size=10.5):
    b = box(s, x, y, w, h, fill=color, radius=True)
    fill_text(b, [[(label, size, txtcolor or BG, True, False)]])
    return b


def arrow(s, x, y, w, h=0.28, color=YELLOW):
    a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_fill(a, color)
    return a


def down_arrow(s, x, y, w=0.35, h=0.4, color=YELLOW):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_fill(a, color)
    return a


def _line_props(conn, color, w, dash=None, tail=False, head=False):
    ln = conn.line
    ln.color.rgb = color; ln.width = Pt(w)
    lnEl = ln._get_or_add_ln()
    if dash is not None:
        lnEl.append(lnEl.makeelement(qn('a:prstDash'), {'val': dash}))
    if head:
        lnEl.append(lnEl.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if tail:
        lnEl.append(lnEl.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))


def connector(s, x1, y1, x2, y2, color=YELLOW, w=1.5, dash=None, tail=True, head=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.shadow.inherit = False
    _line_props(c, color, w, dash, tail, head)
    return c


def fbox(s, shape, x, y, w, h, fill, lines, txtcolor=WHITE, line_color=None, fsize=10.5):
    shp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_fill(shp, fill, line_color, 1 if line_color else None)
    fill_text(shp, [[(l, fsize, txtcolor, True, False)] for l in lines])
    return shp


def seq_msg(s, x_from, x_to, y, label, dashed=False, color=WHITE):
    connector(s, x_from, y, x_to, y, color=color, w=1.25,
              dash=('dash' if dashed else None), tail=True)
    lx = min(x_from, x_to); lw = abs(x_to - x_from)
    txt(s, lx, y - 0.27, lw, 0.24, [[(label, 8.5, GREY, False, False)]], align=PP_ALIGN.CENTER)


def label(s, x, y, w, text, color=YELLOW, size=9):
    txt(s, x, y, w, 0.22, [[(text, size, color, True, False)]], align=PP_ALIGN.CENTER)


def divider(kicker, title, subtitle=""):
    _SEC[0] += 1
    s = slide(BG2)
    accent_bar(s)
    box(s, 0, 7.34, 13.333, 0.16, fill=YELLOW, radius=False)
    txt(s, 1.15, 2.9, 11.0, 0.5, [[("SECTION " + kicker.upper(), 13, YELLOW, True, False)]])
    txt(s, 1.15, 3.3, 11.0, 1.0, [[(title, 38, WHITE, True, False)]])
    box(s, 1.18, 4.55, 3.0, 0.04, fill=YELLOW, radius=False)
    if subtitle:
        txt(s, 1.15, 4.75, 10.8, 0.8, [[(subtitle, 14, GREY, False, False)]], line_spacing=1.1)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.333, 0.16, fill=YELLOW, radius=False)
box(s, 0, 7.34, 13.333, 0.16, fill=YELLOW, radius=False)
txt(s, 0.9, 1.85, 11.5, 0.5, [[("ORDER ASSISTANT AI  ·  FOR KITCHEN & BATH ORDER FULFILLMENT", 14, YELLOW, True, False)]])
txt(s, 0.9, 2.45, 11.6, 1.8,
    [[("Autonomous Purchase Order", 42, WHITE, True, False)],
     [("to Fulfillment Orchestration", 42, WHITE, True, False)]], line_spacing=1.0)
box(s, 0.94, 4.3, 3.0, 0.04, fill=YELLOW, radius=False)
txt(s, 0.9, 4.5, 11.6, 1.4,
    [[("An Order Assistant AI that turns every distributor purchase order the business receives \u2014 ", 15, GREY, False, False)],
     [("from intake to shipment \u2014 into a fast, accurate, self-validating workflow, with a ", 15, GREY, False, False)],
     [("team of specialised AI agents and human-in-the-loop control on exceptions.", 15, GREY, False, False)]],
    line_spacing=1.05)
txt(s, 0.9, 6.5, 11.4, 0.5,
    [[("Business case  \u00b7  Solution architecture  \u00b7  The AI tool we built", 12, DGREY, False, True)]])

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — BUSINESS CONTEXT
# ══════════════════════════════════════════════════════════════════════════════
divider("business context", "Business context",
        "The business need and problem, the solution approach, and the business impact.")

# ── Business need & problem ───────────────────────────────────────────────────
s = slide()
header(s, "Business need & problem", "From distributor PO to shipment \u2014 the manual bottleneck")
band = box(s, 0.6, 1.8, 12.13, 1.0, fill=PANEL2)
box(s, 0.6, 1.8, 0.1, 1.0, fill=YELLOW, radius=False)
txt(s, 0.9, 1.93, 11.6, 0.85,
    [[("The business receives a high volume of purchase orders from wholesale distributors, showrooms and retailers ", 12.5, WHITE, False, False)],
     [("(e.g. Great Lakes Plumbing) across email, EDI and portals. Every order must be validated, priced, sourced "
       "and approved before it can ship \u2014 today that work lands on customer-service reps, order by order, line by line.",
       12.5, GREY, False, False)]], line_spacing=1.05)
probs = [
    ("Constant product churn", "The catalog continually supersedes fixtures, cartridges and finishes. CSRs must catch obsolete SKUs and find the approved replacement on every order."),
    ("Mixed units of measure", "Items sell in each, kits and cases (e.g. Clearflo drain kits). Ordered quantities must be converted correctly or shipments and invoices go wrong."),
    ("Contract & margin complexity", "Distributor-specific contract pricing, promos and rebates \u2014 plus family-level discount / margin policies that must not be breached."),
    ("Account, credit & approval", "Validating the distributor account, the ordering buyer, credit limits and spend approvals \u2014 by hand, across multiple systems."),
]
cw = 5.9; ch = 1.5; gap = 0.35; x0 = 0.6; y0 = 3.05
pos = [(x0, y0), (x0+cw+gap, y0), (x0, y0+ch+0.2), (x0+cw+gap, y0+ch+0.2)]
for (px, py), (t, d) in zip(pos, probs):
    box(s, px, py, cw, ch, fill=PANEL)
    box(s, px, py, 0.09, ch, fill=RED, radius=False)
    txt(s, px+0.32, py+0.18, cw-0.55, 0.4, [[(t, 15, WHITE, True, False)]])
    txt(s, px+0.32, py+0.62, cw-0.55, 0.8, [[(d, 11.5, GREY, False, False)]], line_spacing=1.03)
imp = box(s, 0.6, 6.55, 12.13, 0.5, fill=BG, line_color=RED, line_w=1)
txt(s, 0.85, 6.63, 11.6, 0.4,
    [[("Impact of the status quo:  ", 12, RED, True, False),
      ("slow order turnaround, wrong-SKU shipments and returns, margin leakage, CSRs buried in lookups.", 12, GREY, False, False)]])

# ── Business solution ─────────────────────────────────────────────────────────
s = slide()
header(s, "Business solution", "An Order Assistant AI with human-in-the-loop control")
txt(s, 0.6, 1.7, 12.1, 0.6,
    [[("An orchestration engine runs a pipeline of specialist AI agents. Each agent makes master-data-driven decisions "
       "automatically and only pauses for a human when a decision is ambiguous or breaches policy.", 13.5, GREY, False, False)]],
    line_spacing=1.1)
cards = [
    ("Straight-through", "Clean orders are validated, priced, sourced, planned and created with zero human touch.", GREEN),
    ("Human-in-the-loop", "Risky or ambiguous decisions pause the pipeline and hand off to a CSR, then resume automatically.", AMBER),
    ("Master-data driven", "Decisions come from governed master data \u2014 not hard-coded rules \u2014 so behaviour is transparent and tunable.", YELLOW),
    ("Fully auditable", "Every automated and human decision is captured with its rationale for governance and traceability.", BLUE),
]
cw = 2.92; gap = 0.13; x0 = 0.6; y0 = 2.65; chh = 3.7
for i, (t, d, col) in enumerate(cards):
    px = x0 + i*(cw+gap)
    box(s, px, y0, cw, chh, fill=PANEL)
    box(s, px, y0, cw, 0.12, fill=col, radius=False)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px+0.32), Inches(y0+0.42), Inches(0.55), Inches(0.55))
    _set_fill(circ, col)
    fill_text(circ, [[(str(i+1), 18, BG, True, False)]])
    txt(s, px+0.3, y0+1.2, cw-0.55, 0.7, [[(t, 16, WHITE, True, False)]])
    txt(s, px+0.3, y0+1.95, cw-0.55, 1.6, [[(d, 12, GREY, False, False)]], line_spacing=1.1)

# ── Business impact ───────────────────────────────────────────────────────────
s = slide()
header(s, "Business impact", "How the Order Assistant AI improves the business")
bx = box(s, 0.6, 1.85, 6.0, 3.45, fill=PANEL)
box(s, 0.6, 1.85, 6.0, 0.12, fill=RED, radius=False)
txt(s, 0.9, 2.07, 5.5, 0.4, [[("TODAY  \u2014  MANUAL", 11, RED, True, False)]])
for i, t in enumerate(["Minutes of multi-system lookups per line",
                       "Obsolete SKUs slip through \u2192 returns",
                       "Inconsistent discounting \u2192 margin leak",
                       "Hard to scale at peak season",
                       "Limited audit of who decided what"]):
    yy = 2.55 + i*0.52
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(yy+0.05), Inches(0.14), Inches(0.14)); _set_fill(d, RED)
    txt(s, 1.22, yy, 5.1, 0.5, [[(t, 12, GREY, False, False)]])
ax = box(s, 6.73, 1.85, 6.0, 3.45, fill=PANEL)
box(s, 6.73, 1.85, 6.0, 0.12, fill=GREEN, radius=False)
txt(s, 7.03, 2.07, 5.5, 0.4, [[("WITH ORDER ASSISTANT AI", 11, GREEN, True, False)]])
for i, t in enumerate(["Clean orders auto-processed end-to-end",
                       "Obsolete SKUs auto-flagged with approved substitute",
                       "Margin policy enforced on every line",
                       "Scales elastically with order volume",
                       "Every decision captured for audit"]):
    yy = 2.55 + i*0.52
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.08), Inches(yy+0.05), Inches(0.14), Inches(0.14)); _set_fill(d, GREEN)
    txt(s, 7.35, yy, 5.1, 0.5, [[(t, 12, GREY, False, False)]])
txt(s, 0.6, 5.55, 12.0, 0.4, [[("Business outcomes", 15, WHITE, True, False)]])
outs = ["Faster distributor turnaround", "Fewer returns & wrong shipments", "Protected margins",
        "Higher CSR productivity", "Consistent, governed decisions"]
cw = 2.35; gap = 0.12; x0 = 0.6; y0 = 6.05
for i, o in enumerate(outs):
    b = box(s, x0+i*(cw+gap), y0, cw, 0.8, fill=PANEL2)
    box(s, x0+i*(cw+gap), y0, cw, 0.08, fill=YELLOW, radius=False)
    fill_text(b, [[(o, 10.5, WHITE, True, False)]])

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 2 — SOLUTION ARCHITECTURE & DESIGN
# ══════════════════════════════════════════════════════════════════════════════
divider("solution architecture & design", "Solution architecture & design",
        "The reference architecture, the runtime component & data-flow view, the process flow, and the request sequence.")

# ── Flagship reference architecture ───────────────────────────────────────────
s = slide()
header(s, "Reference architecture", "Solution architecture \u2014 layered reference view")

LX, LW = 0.6, 1.68          # tier-label column
CX = 2.42                   # content start
CW = 8.02                   # content width (to 10.44)
XR, XRW = 10.66, 2.07       # cross-cutting column (to 12.73)


def tier_label(y, h, tag, color):
    b = box(s, LX, y, LW, h, fill=PANEL)
    box(s, LX, y, 0.08, h, fill=color, radius=False)
    fill_text(b, [[(w, 9.5, WHITE, True, False)] for w in tag.split("\n")])


def tiles(y, h, items, fill=PANEL2, line=LINE, tcol=WHITE, fs=9.5):
    n = len(items); g = 0.14
    w = (CW - g*(n-1)) / n
    for i, it in enumerate(items):
        px = CX + i*(w+g)
        b = box(s, px, y, w, h, fill=fill, line_color=line, line_w=0.75)
        fill_text(b, [[(ln, fs, tcol, True, False)] for ln in it.split("\n")])


# tiers
tier_label(1.85, 0.6, "CHANNELS &\nINGESTION", BLUE)
tiles(1.85, 0.6, ["Email / PDF", "EDI", "Web portal", "Manual entry"])

tier_label(2.57, 0.62, "EXPERIENCE\nLAYER", TEAL)
eb = box(s, CX, 2.57, CW, 0.62, fill=PANEL2, line_color=YELLOW, line_w=1)
fill_text(eb, [[("CSR Workspace  \u00b7  Streamlit", 11, WHITE, True, False)],
               [("PO intake  \u00b7  live agent panels  \u00b7  one-click decisions  \u00b7  audit viewer", 8.5, GREY, False, False)]])

tier_label(3.31, 1.16, "AI\nORCHESTRATION", YELLOW)
ob = box(s, CX, 3.31, CW, 1.16, fill=PANEL)
txt(s, CX+0.15, 3.37, CW-0.3, 0.3, [[("Orchestrator  \u2014  resumable pipeline / state machine  (pause on exception \u00b7 resume on decision)", 9.5, YELLOW, True, False)]])
agents = ["Intake", "Customer\nValidation", "Product\nMatching", "Pricing &\nPromo", "Credit",
          "Inventory", "Shipments", "Approvals", "Order\nExecution"]
n = len(agents); g = 0.08; aw = (CW-0.3 - g*(n-1))/n
for i, a in enumerate(agents):
    px = CX+0.15 + i*(aw+g)
    b = box(s, px, 3.72, aw, 0.65, fill=PANEL2, line_color=YELLOW, line_w=0.75)
    fill_text(b, [[(ln, 8, WHITE, True, False)] for ln in a.split("\n")])

tier_label(4.59, 0.62, "DOMAIN\nSERVICES", PURPLE)
tiles(4.59, 0.62, ["Pricing engine", "Credit service", "Inventory / ATP", "Logistics", "Compliance", "Budget / approval"], fs=9)

tier_label(5.33, 1.02, "DATA &\nINTEGRATION", GREEN)
mdp = box(s, CX, 5.33, 5.15, 1.02, fill=PANEL2)
txt(s, CX+0.12, 5.38, 4.9, 0.25, [[("MASTER-DATA FOUNDATION", 8.5, YELLOW, True, False)]])
md = ["Customer", "Buyer", "Product", "Pricing", "Credit", "Inventory", "Logistics", "Budget/Comp."]
mg = 0.06; mw = (5.15-0.24 - mg*3)/4
for i, m in enumerate(md):
    r, c = divmod(i, 4)
    px = CX+0.12 + c*(mw+mg); py = 5.63 + r*0.34
    b = box(s, px, py, mw, 0.3, fill=BG, line_color=LINE, line_w=0.5)
    fill_text(b, [[(m, 7.5, GREY, True, False)]])
intp = box(s, CX+5.32, 5.33, CW-5.32, 1.02, fill=PANEL2)
txt(s, CX+5.44, 5.38, CW-5.5, 0.25, [[("INTEGRATION ADAPTERS", 8.5, YELLOW, True, False)]])
for i, t in enumerate(["ERP / Order Mgmt", "Email / SMTP", "Audit store"]):
    b = box(s, CX+5.44, 5.63 + i*0.23, CW-5.56, 0.2, fill=BG, line_color=LINE, line_w=0.5)
    fill_text(b, [[(t, 7.8, GREY, True, False)]])

# cross-cutting column
cc = box(s, XR, 1.85, XRW, 4.5, fill=PANEL)
box(s, XR, 1.85, XRW, 0.1, fill=YELLOW, radius=False)
txt(s, XR+0.14, 2.0, XRW-0.28, 0.3, [[("CROSS-CUTTING", 9, YELLOW, True, False)]])
for i, (t, sub) in enumerate([("Exception\nGovernance & routing", ""),
                              ("Audit &\nTraceability", ""),
                              ("Security &\nAccess control", ""),
                              ("Observability\n& Logging", "")]):
    yy = 2.34 + i*0.98
    b = box(s, XR+0.14, yy, XRW-0.28, 0.85, fill=PANEL2, line_color=LINE, line_w=0.75)
    fill_text(b, [[(ln, 9, WHITE, True, False)] for ln in t.split("\n")])

# top-down flow arrows on far left
for gy in (2.47, 3.20, 4.48, 5.22):
    down_arrow(s, 1.29, gy, 0.24, 0.1, color=YELLOW)
txt(s, 0.6, 6.5, 12.1, 0.35,
    [[("Logical reference architecture.  ", 9.5, YELLOW, True, True),
      ("POC runs on Streamlit + Python with master data as workbooks; production maps to enterprise ERP / MDM, identity and cloud services.",
       9.5, DGREY, False, True)]])

# ── Component & data-flow view ────────────────────────────────────────────────
s = slide()
header(s, "Runtime view", "Component & data-flow architecture")
src = box(s, 0.6, 1.95, 2.55, 1.35, fill=PANEL2)
txt(s, 0.78, 2.05, 2.2, 0.3, [[("PO SOURCES", 10, YELLOW, True, False)]])
for i, t in enumerate(["Email / PDF", "EDI", "Portal / Manual"]):
    chip(s, 0.78, 2.42 + i*0.28, 2.2, 0.24, t, BG, GREY, size=9)
ui = box(s, 3.55, 1.95, 6.35, 0.95, fill=PANEL2, line_color=YELLOW, line_w=1)
txt(s, 3.75, 2.05, 6.0, 0.3, [[("EXPERIENCE LAYER", 10, YELLOW, True, False)]])
txt(s, 3.75, 2.34, 6.0, 0.5, [[("CSR Workspace  \u2014  Streamlit UI  (intake \u00b7 live agent panels \u00b7 one-click decisions \u00b7 audit)", 11, WHITE, True, False)]], line_spacing=0.95)
integ = box(s, 10.25, 1.95, 2.45, 3.95, fill=PANEL2)
txt(s, 10.43, 2.05, 2.1, 0.3, [[("INTEGRATIONS", 10, YELLOW, True, False)]])
for i, t in enumerate(["ERP / Order Mgmt", "Email / SMTP", "Approval routing", "Audit store"]):
    chip(s, 10.43, 2.45 + i*0.55, 2.1, 0.42, t, BG, GREY, size=9.5)
eng = box(s, 3.55, 3.2, 6.35, 2.7, fill=PANEL)
box(s, 3.55, 3.2, 6.35, 0.1, fill=YELLOW, radius=False)
txt(s, 3.75, 3.34, 6.0, 0.3, [[("ORCHESTRATION ENGINE", 10.5, YELLOW, True, False),
    ("   (resumable pipeline / state machine)", 9, DGREY, False, True)]])
comp = [("PO Extractor", 3.72, 3.75), ("Intake Resolver", 5.87, 3.75), ("Account Validator", 8.02, 3.75),
        ("Decision Pipeline\n(9 agents)", 3.72, 4.5), ("Exception\nGovernance", 5.87, 4.5), ("Order\nExecution", 8.02, 4.5)]
for (t, px, py) in comp:
    b = box(s, px, py, 1.9, 0.65 if "\n" not in t else 1.05, fill=PANEL2, line_color=LINE, line_w=0.75)
    fill_text(b, [[(ln, 9.5, WHITE, True, False)] for ln in t.split("\n")])
md = box(s, 3.55, 6.15, 6.35, 0.75, fill=PANEL2)
txt(s, 3.75, 6.22, 6.0, 0.3, [[("MASTER-DATA FOUNDATION", 10, YELLOW, True, False)]])
for i, m in enumerate(["Customer", "Buyer", "Product", "Pricing", "Credit", "Inventory", "Logistics", "Budget"]):
    mb = box(s, 3.75 + i*0.755, 6.5, 0.72, 0.32, fill=BG, line_color=LINE, line_w=0.5)
    fill_text(mb, [[(m, 7.5, GREY, True, False)]])
connector(s, 3.15, 2.62, 3.55, 2.62, color=YELLOW, w=1.5, tail=True)
connector(s, 6.7, 2.9, 6.7, 3.2, color=YELLOW, w=1.5, tail=True, head=True)
connector(s, 6.7, 5.9, 6.7, 6.15, color=YELLOW, w=1.5, tail=True, head=True)
connector(s, 9.9, 4.3, 10.25, 4.3, color=YELLOW, w=1.5, tail=True, head=True)
txt(s, 0.6, 7.02, 9.5, 0.3, [[("Data-driven decisions read governed master data; results & exceptions flow back to the CSR and out to ERP / email.", 9, DGREY, False, True)]])

# ── Process flowchart ─────────────────────────────────────────────────────────
s = slide()
header(s, "Process design", "Process flowchart")
FT = MSO_SHAPE.FLOWCHART_TERMINATOR
FP = MSO_SHAPE.FLOWCHART_PROCESS
FD = MSO_SHAPE.FLOWCHART_DECISION
cx = 2.35; cw = 3.9; cxc = cx + cw/2
fbox(s, FT, cx, 1.8, cw, 0.5, GREEN, ["PO received"], BG)
fbox(s, FP, cx, 2.42, cw, 0.6, PANEL2, ["Intake \u2014 extract & reconcile vs master data"], WHITE, LINE)
fbox(s, FD, cx+0.45, 3.12, cw-0.9, 0.85, PANEL, ["Intake issues?"], WHITE, BLUE)
fbox(s, FP, cx, 4.08, cw, 0.6, PANEL2, ["Run decision pipeline (Customer \u2192 \u2026 \u2192 Approvals)"], WHITE, LINE)
fbox(s, FD, cx+0.45, 4.78, cw-0.9, 0.85, PANEL, ["Policy exception?"], WHITE, AMBER)
fbox(s, FP, cx, 5.74, cw, 0.55, PANEL2, ["Order Execution \u2014 create sales order"], WHITE, LINE)
fbox(s, FT, cx, 6.4, cw, 0.5, GREEN, ["Order created + full audit trail"], BG)
rx = 7.9; rw = 4.2
fbox(s, FP, rx, 3.02, rw, 0.9, PANEL, ["CSR resolves gate(s):", "substitute \u00b7 qty \u00b7 UOM \u00b7 buyer"], WHITE, BLUE, 10)
fbox(s, FP, rx, 4.68, rw, 0.9, PANEL, ["CSR decides:", "Approve \u00b7 Reject \u00b7 Notify"], WHITE, AMBER, 10)
fbox(s, FT, rx, 5.9, rw, 0.55, RED, ["Order stopped / routed to team"], WHITE)
connector(s, cxc, 2.3, cxc, 2.42, color=YELLOW, tail=True)
connector(s, cxc, 3.02, cxc, 3.12, color=YELLOW, tail=True)
connector(s, cxc, 3.97, cxc, 4.08, color=YELLOW, tail=True); label(s, cxc+0.08, 3.99, 0.6, "No", GREY, 8.5)
connector(s, cxc, 4.68, cxc, 4.78, color=YELLOW, tail=True)
connector(s, cxc, 5.63, cxc, 5.74, color=YELLOW, tail=True); label(s, cxc+0.08, 5.64, 0.6, "No", GREY, 8.5)
connector(s, cxc, 6.29, cxc, 6.4, color=YELLOW, tail=True)
connector(s, cx+cw-0.45, 3.545, rx, 3.47, color=BLUE, tail=True); label(s, cx+cw-0.3, 3.2, 0.7, "Yes", BLUE, 8.5)
connector(s, rx, 3.9, rx-0.35, 3.9, color=BLUE, tail=False)
connector(s, rx-0.35, 3.9, rx-0.35, 4.38, color=BLUE, tail=False)
connector(s, rx-0.35, 4.38, cxc, 4.38, color=BLUE, tail=True); label(s, cxc+1.6, 4.14, 1.6, "resolved \u2192 continue", BLUE, 8)
connector(s, cx+cw-0.45, 5.205, rx, 5.13, color=AMBER, tail=True); label(s, cx+cw-0.3, 4.86, 0.7, "Yes", AMBER, 8.5)
connector(s, rx, 5.58, rx-0.35, 5.58, color=AMBER, tail=False)
connector(s, rx-0.35, 5.58, rx-0.35, 6.015, color=AMBER, tail=False)
connector(s, rx-0.35, 6.015, cx+cw, 6.015, color=AMBER, tail=True); label(s, cx+cw+0.15, 5.78, 1.5, "approve \u2192 continue", AMBER, 8)
connector(s, rx+rw/2, 5.58, rx+rw/2, 5.9, color=RED, tail=True); label(s, rx+rw/2+0.1, 5.6, 0.9, "reject", RED, 8.5)

# ── Sequence diagram ──────────────────────────────────────────────────────────
s = slide()
header(s, "Interaction design", "Sequence diagram")
parts = [("CSR / Buyer", 1.5), ("CSR Workspace UI", 3.55), ("Orchestrator", 5.6),
         ("Specialist Agents", 7.65), ("Master Data", 9.7), ("ERP / Email", 11.75)]
top = 1.95; bottom = 6.75
for (name, cxp) in parts:
    hb = box(s, cxp-0.95, top, 1.9, 0.5, fill=PANEL2, line_color=YELLOW, line_w=0.75)
    fill_text(hb, [[(name, 10, WHITE, True, False)]])
    connector(s, cxp, top+0.5, cxp, bottom, color=LINE, w=1.25, dash='dash', tail=False)
C = {name: cxp for name, cxp in parts}
seq_msg(s, C["CSR / Buyer"], C["CSR Workspace UI"], 2.85, "1. Submit PO")
seq_msg(s, C["CSR Workspace UI"], C["Orchestrator"], 3.2, "2. start(po)")
seq_msg(s, C["Orchestrator"], C["Specialist Agents"], 3.55, "3. run intake + pipeline")
seq_msg(s, C["Specialist Agents"], C["Master Data"], 3.9, "4. lookup")
seq_msg(s, C["Master Data"], C["Specialist Agents"], 4.25, "5. records", dashed=True, color=GREY)
seq_msg(s, C["Specialist Agents"], C["Orchestrator"], 4.6, "6. results / exceptions", dashed=True, color=GREY)
frx = C["CSR Workspace UI"] - 1.15
frw = (C["Orchestrator"] - C["CSR Workspace UI"]) + 2.3
fr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(frx), Inches(4.85), Inches(frw), Inches(1.15))
fr.fill.background(); fr.line.color.rgb = AMBER; fr.line.width = Pt(1); fr.shadow.inherit = False
tab = box(s, frx, 4.85, 1.5, 0.28, fill=AMBER, radius=False)
fill_text(tab, [[("alt  [exception]", 8.5, BG, True, False)]])
seq_msg(s, C["Orchestrator"], C["CSR Workspace UI"], 5.3, "7. CSR DECISION NEEDED", color=AMBER)
seq_msg(s, C["CSR / Buyer"], C["CSR Workspace UI"], 5.62, "8. decision (pick / enter / approve)")
seq_msg(s, C["CSR Workspace UI"], C["Orchestrator"], 5.92, "9. resume(decision)")
seq_msg(s, C["Orchestrator"], C["ERP / Email"], 6.3, "10. create order / notify")
seq_msg(s, C["ERP / Email"], C["Orchestrator"], 6.6, "11. order # / ack", dashed=True, color=GREY)
txt(s, 0.6, 7.0, 12.0, 0.3, [[("The exception block (7\u20139) repeats per gate; clean orders skip it entirely (straight-through).", 9, DGREY, False, True)]])

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 3 — THE AI TOOL WE BUILT
# ══════════════════════════════════════════════════════════════════════════════
divider("the AI tool we built", "The AI tool we built",
        "The agents, the master-data foundation, the pipeline, the two processing paths, governance and technology.")

# ── The AI agent team ─────────────────────────────────────────────────────────
s = slide()
header(s, "What we built", "The AI agent team")
team = [
    ("Intake Agent", "Digitises the raw PO, extracts header/buyer/ship-to/lines, and reconciles every line against master data.", "Obsolete SKU \u00b7 bad qty \u00b7 UOM \u00b7 unknown buyer"),
    ("Customer Validation Agent", "Resolves customer & account hierarchy, validates ship-to, and authorises the buyer.", "Account \u00b7 ship-to \u00b7 buyer authority"),
    ("Product Matching Agent", "Confirms each SKU against the catalog and runs the compliance sub-check.", "Catalog match \u00b7 compliance"),
    ("Pricing & Promo Agent", "Builds the price waterfall (list \u2192 contract \u2192 tier \u2192 promo \u2192 rebate \u2192 net) and totals.", "Margin / discount policy"),
    ("Credit Agent", "Checks order value against available credit and payment terms.", "Credit limit \u00b7 terms \u00b7 risk"),
    ("Inventory Agent", "Sources each line across distribution centres and confirms availability.", "ATP \u00b7 allocation"),
    ("Shipments Agent", "Confirms carrier serviceability, freight and delivery SLA.", "Fulfillment plan \u00b7 ETA"),
    ("Approvals Agent", "Runs the approval-matrix / budget check last, before execution.", "Budget \u00b7 self-approval"),
    ("Order Execution Agent", "Creates the sales order and emits the completion record.", "Order created \u00b7 audit"),
]
cw = 3.95; ch = 1.55; gapx = 0.13; gapy = 0.16; x0 = 0.6; y0 = 1.85
for i, (t, d, tag) in enumerate(team):
    r, c = divmod(i, 3)
    px = x0 + c*(cw+gapx); py = y0 + r*(ch+gapy)
    box(s, px, py, cw, ch, fill=PANEL)
    box(s, px, py, 0.08, ch, fill=YELLOW, radius=False)
    txt(s, px+0.28, py+0.16, cw-0.5, 0.4, [[("\U0001F916  " + t, 12.5, WHITE, True, False)]])
    txt(s, px+0.28, py+0.62, cw-0.5, 0.7, [[(d, 10.5, GREY, False, False)]], line_spacing=1.02)
    txt(s, px+0.28, py+1.24, cw-0.5, 0.3, [[(tag, 9, YELLOW, False, True)]])

# ── Master data foundation ────────────────────────────────────────────────────
s = slide()
header(s, "What we built", "Master data drives every decision")
txt(s, 0.6, 1.75, 12.0, 0.5,
    [[("Each agent reads governed master data (mock Excel workbooks in the POC; enterprise systems in production). "
       "Change the data, not the code, to change behaviour.", 13, GREY, False, False)]], line_spacing=1.1)
data = [
    ("Customer", "Account hierarchy, tiers, buying history, fulfillment rules"),
    ("Buyer", "Buyer profiles, roles, permissions, product-family access"),
    ("Product", "Catalog, lifecycle status, substitutes, UOM conversions"),
    ("Pricing", "Price lists, contracts, volume tiers, promos, rebates"),
    ("Credit", "Credit limits, available credit, payment terms, risk"),
    ("Inventory", "Plant / DC stock, in-transit, ATP, allocation"),
    ("Logistics", "Carriers, serviceability, freight, delivery SLA"),
    ("Budget", "Cost-centre budgets, approval matrix, buyer authority"),
    ("Compliance", "Regional restrictions, approvals, required documents"),
    ("Governance", "Exception routing, escalation targets, SLAs"),
]
cw = 3.95; ch = 0.92; gapx = 0.13; gapy = 0.14; x0 = 0.6; y0 = 2.55
for i, (t, d) in enumerate(data):
    r, c = divmod(i, 3)
    px = x0 + c*(cw+gapx); py = y0 + r*(ch+gapy)
    box(s, px, py, cw, ch, fill=PANEL2)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px+0.22), Inches(py+0.3), Inches(0.28), Inches(0.28)); _set_fill(dot, YELLOW)
    txt(s, px+0.66, py+0.14, cw-0.85, 0.35, [[(t, 12.5, WHITE, True, False)]])
    txt(s, px+0.66, py+0.5, cw-0.85, 0.4, [[(d, 9.5, GREY, False, False)]], line_spacing=1.0)

# ── Orchestration pipeline ────────────────────────────────────────────────────
s = slide()
header(s, "What we built", "The orchestration pipeline")
txt(s, 0.6, 1.7, 12.0, 0.45,
    [[("Sequential decision layers thread a shared context. The pipeline pauses on the first exception and resumes after the CSR decides.", 12.5, GREY, False, False)]], line_spacing=1.05)
stages = ["Intake", "Customer\nValidation", "Product\nMatching", "Pricing &\nPromo", "Credit", "Inventory", "Shipments", "Approvals"]
n = len(stages); bw = 1.28; bh = 1.0; gap = 0.145; x0 = 0.62; y0 = 2.55
for i, stg in enumerate(stages):
    px = x0 + i*(bw+gap)
    b = box(s, px, y0, bw, bh, fill=PANEL, line_color=YELLOW, line_w=1)
    fill_text(b, [[(str(i+1), 10, YELLOW, True, False)]] + [[(ln, 10, WHITE, True, False)] for ln in stg.split("\n")])
    if i < n-1:
        arrow(s, px+bw-0.02, y0+bh/2-0.13, gap+0.06, 0.26)
down_arrow(s, x0 + (n-1)*(bw+gap) + bw/2 - 0.2, y0+bh+0.05, 0.4, 0.32)
gbox = box(s, 8.6, 3.95, 4.12, 0.72, fill=PANEL2, line_color=GREEN, line_w=1)
fill_text(gbox, [[("Order Execution Agent  \u00b7  create order", 11.5, WHITE, True, False)]])
govb = box(s, 0.62, 3.95, 7.7, 0.72, fill=PANEL2)
fill_text(govb, [[("Exception Governance & Routing  \u2014  runs on every order (records + routes exceptions)", 11, GREY, True, False)]])
notes = [
    ("Pauses on first exception", "The pipeline stops at the layer that needs a decision \u2014 nothing downstream runs on bad data.", AMBER),
    ("Resumes automatically", "Once the CSR decides, the shared context is updated and the pipeline continues from where it paused.", GREEN),
    ("Approvals run last", "Budget / approval-matrix check runs only after product, pricing, credit, inventory and logistics pass.", YELLOW),
]
cw = 3.95; x0n = 0.6; y0n = 5.15; chh = 1.5
for i, (t, d, col) in enumerate(notes):
    px = x0n + i*(cw+0.13)
    box(s, px, y0n, cw, chh, fill=PANEL)
    box(s, px, y0n, 0.08, chh, fill=col, radius=False)
    txt(s, px+0.28, y0n+0.2, cw-0.5, 0.4, [[(t, 13, WHITE, True, False)]])
    txt(s, px+0.28, y0n+0.68, cw-0.5, 0.7, [[(d, 11, GREY, False, False)]], line_spacing=1.05)

# ── One engine, two paths ─────────────────────────────────────────────────────
s = slide()
header(s, "What we built", "One engine, two paths")
box(s, 0.6, 1.95, 5.95, 4.6, fill=PANEL)
box(s, 0.6, 1.95, 5.95, 0.13, fill=GREEN, radius=False)
txt(s, 0.9, 2.25, 5.4, 0.5, [[("HAPPY FLOW", 12, GREEN, True, False)]])
txt(s, 0.9, 2.6, 5.4, 0.5, [[("Straight-through processing", 18, WHITE, True, False)]])
for i, t in enumerate(["Clean PO, no ambiguity or policy breach", "Every agent validates & passes automatically",
                       "Zero human intervention", "Order created in a single pass", "Full audit trail produced"]):
    yy = 3.25 + i*0.55
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(yy+0.06), Inches(0.16), Inches(0.16)); _set_fill(dot, GREEN)
    txt(s, 1.25, yy, 5.1, 0.5, [[(t, 12.5, GREY, False, False)]])
chip(s, 0.9, 6.05, 5.3, 0.42, "Outcome:  order created with no CSR touch", PANEL2, WHITE)
box(s, 6.78, 1.95, 5.95, 4.6, fill=PANEL)
box(s, 6.78, 1.95, 5.95, 0.13, fill=AMBER, radius=False)
txt(s, 7.08, 2.25, 5.4, 0.5, [[("CSR APPROVAL FLOW", 12, AMBER, True, False)]])
txt(s, 7.08, 2.6, 5.4, 0.5, [[("Human-in-the-loop", 18, WHITE, True, False)]])
for i, t in enumerate(["AI auto-resolves everything it can from data", "Ambiguous / policy-breaching items PAUSE",
                       "CSR gets a pre-investigated, one-click decision", "Guardrails validate every manual entry",
                       "Pipeline resumes automatically after each gate"]):
    yy = 3.25 + i*0.55
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.13), Inches(yy+0.06), Inches(0.16), Inches(0.16)); _set_fill(dot, AMBER)
    txt(s, 7.43, yy, 5.1, 0.5, [[(t, 12.5, GREY, False, False)]])
chip(s, 7.08, 6.05, 5.3, 0.42, "Outcome:  CSR decides only the exceptions", PANEL2, WHITE)

# ── Happy Flow detail ─────────────────────────────────────────────────────────
s = slide()
header(s, "Demo 1", "Happy Flow \u2014 straight-through processing")
steps = ["Submit\nPO", "Intake\nextract + review", "Customer\nValidation", "Product\nMatching", "Pricing",
         "Credit", "Inventory", "Shipments", "Approvals", "Order\nCreated"]
bw = 1.12; bh = 0.95; gap = 0.055; x0 = 0.62; y0 = 2.0
for i, stp in enumerate(steps):
    px = x0 + i*(bw+gap); last = (i == len(steps)-1)
    b = box(s, px, y0, bw, bh, fill=(GREEN if last else PANEL), line_color=(None if last else GREEN), line_w=0.75)
    col = BG if last else WHITE
    fill_text(b, [[(ln, 9, col, True, False)] for ln in stp.split("\n")])
    if i < len(steps)-1:
        arrow(s, px+bw-0.02, y0+bh/2-0.1, gap+0.05, 0.2, color=GREEN)
txt(s, 0.62, 3.1, 12.0, 0.3, [[("All checks pass automatically \u2014 no ", 12, GREY, False, False),
    ("CSR DECISION NEEDED", 12, GREEN, True, False), (" banners appear.", 12, GREY, False, False)]])
box(s, 0.6, 3.6, 6.5, 2.95, fill=PANEL)
txt(s, 0.85, 3.78, 6.0, 0.4, [[("EXAMPLE PO  \u2014  PO-2026-30001", 11, YELLOW, True, False)]])
lines = ["Great Lakes Plumbing Supply Co  \u00b7  john.miller@glps.com", "Ship to: Chicago DC, IL 60639", "",
         "1  SKU-CTG-4520  Ceramic Disc Cartridge     100 EA", "2  SKU-SEL-1150  Tank-to-Bowl Gasket Kit    120 EA",
         "3  SKU-VLV-2201  Pressure-Balancing Valve    15 EA"]
txt(s, 0.85, 4.2, 6.0, 2.2, [[(l, 11, GREY if not l.startswith(("1", "2", "3")) else WHITE, False, False)] for l in lines], line_spacing=1.15)
box(s, 7.28, 3.6, 5.45, 2.95, fill=PANEL2)
box(s, 7.28, 3.6, 0.1, 2.95, fill=GREEN, radius=False)
txt(s, 7.55, 3.78, 5.0, 0.4, [[("AUTOMATED RESULT", 11, GREEN, True, False)]])
for i, (k, v) in enumerate([("Subtotal", "$3,083.00"), ("Freight", "$151.62"), ("Tax", "$271.79")]):
    yy = 4.25 + i*0.5
    txt(s, 7.55, yy, 3.0, 0.4, [[(k, 12.5, GREY, False, False)]])
    txt(s, 10.3, yy, 2.2, 0.4, [[(v, 12.5, WHITE, True, False)]], align=PP_ALIGN.RIGHT)
box(s, 7.55, 5.75, 4.95, 0.02, fill=LINE, radius=False)
txt(s, 7.55, 5.85, 3.0, 0.4, [[("ORDER TOTAL", 14, YELLOW, True, False)]])
txt(s, 10.0, 5.85, 2.5, 0.4, [[("$3,529.53", 15, YELLOW, True, False)]], align=PP_ALIGN.RIGHT)

# ── CSR Approval Flow detail ──────────────────────────────────────────────────
s = slide()
header(s, "Demo 2", "CSR Approval Flow \u2014 five decision gates")
txt(s, 0.6, 1.72, 12.1, 0.4,
    [[("A single PO engineered to trigger five gates. The AI investigates, presents a one-click choice, captures the decision, and resumes.", 12, GREY, False, False)]], line_spacing=1.05)
gates = [
    ("1", "Obsolete product", "Line 1 SKU is obsolete \u2014 AI proposes approved substitutes with price impact.", "INTAKE"),
    ("2", "Invalid quantity", "Line 4 quantity is 0 \u2014 CSR enters a valid, positive whole number (validated).", "INTAKE"),
    ("3", "UOM conversion", "Line 5 ordered in EA but sold in KIT \u2014 CSR picks the kit quantity to fulfil.", "INTAKE"),
    ("4", "Unresolved buyer", "PO email isn't a registered buyer \u2014 CSR selects from the account's buyers.", "INTAKE"),
    ("5", "Pricing / margin exception", "Discount 21.5% exceeds the 10% family policy \u2014 CSR approves the override.", "PIPELINE"),
]
cw = 2.35; ch = 3.0; gap = 0.12; x0 = 0.6; y0 = 2.35
for i, (num, t, d, phase) in enumerate(gates):
    px = x0 + i*(cw+gap); pcol = BLUE if phase == "INTAKE" else AMBER
    box(s, px, y0, cw, ch, fill=PANEL)
    box(s, px, y0, cw, 0.1, fill=pcol, radius=False)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px+0.28), Inches(y0+0.3), Inches(0.55), Inches(0.55)); _set_fill(circ, pcol)
    fill_text(circ, [[(num, 18, BG, True, False)]])
    txt(s, px+0.24, y0+1.0, cw-0.45, 0.7, [[(t, 13, WHITE, True, False)]], line_spacing=0.95)
    txt(s, px+0.24, y0+1.75, cw-0.45, 1.0, [[(d, 10, GREY, False, False)]], line_spacing=1.05)
    chip(s, px+0.24, y0+2.62, cw-0.48, 0.32, phase, pcol, BG)
txt(s, 0.6, 5.6, 12.0, 0.4, [[("Blue = INTAKE gates (data reconciliation)      Amber = PIPELINE gate (policy breach)", 11, GREY, False, True)]])
box(s, 0.6, 6.0, 12.13, 0.62, fill=PANEL2)
box(s, 0.6, 6.0, 0.1, 0.62, fill=YELLOW, radius=False)
txt(s, 0.9, 6.12, 11.6, 0.4,
    [[("Approx. order total after CSR choices: ", 12, GREY, False, False), ("$10,659.50", 13, YELLOW, True, False),
      ("   \u2014   five multi-system lookups reduced to five one-click decisions.", 12, GREY, False, False)]])

# ── Exception model & routing ─────────────────────────────────────────────────
s = slide()
header(s, "What we built", "Exception model & escalation routing")
box(s, 0.6, 1.9, 6.0, 2.3, fill=PANEL)
box(s, 0.6, 1.9, 0.1, 2.3, fill=BLUE, radius=False)
txt(s, 0.9, 2.08, 5.5, 0.4, [[("INTAKE GATES  \u2014  data reconciliation", 12.5, BLUE, True, False)]])
txt(s, 0.9, 2.55, 5.5, 1.5,
    [[("Obsolete / substitute SKU  \u00b7  unresolved or missing SKU", 11.5, GREY, False, False)],
     [("Invalid quantity  \u00b7  unit-of-measure conversion", 11.5, GREY, False, False)],
     [("Unresolved buyer  \u00b7  partial / unresolved ship-to", 11.5, GREY, False, False)]], line_spacing=1.2)
box(s, 6.73, 1.9, 6.0, 2.3, fill=PANEL)
box(s, 6.73, 1.9, 0.1, 2.3, fill=AMBER, radius=False)
txt(s, 7.03, 2.08, 5.5, 0.4, [[("PIPELINE GATES  \u2014  policy breach", 12.5, AMBER, True, False)]])
txt(s, 7.03, 2.55, 5.5, 1.5,
    [[("Pricing / margin-policy exception", 11.5, GREY, False, False)],
     [("Credit hold  \u00b7  budget exceeded / approval required", 11.5, GREY, False, False)],
     [("Compliance restriction  \u00b7  inventory / logistics constraints", 11.5, GREY, False, False)]], line_spacing=1.2)
txt(s, 0.6, 4.45, 12.0, 0.4, [[("Every exception is routed to the responsible team", 15, WHITE, True, False)]])
routes = [("Product Specialist", "SKU / substitution"), ("Order Ops Supervisor", "Quantity / intake"),
          ("Sales Ops / Account Mgr", "Buyer / ship-to"), ("Pricing Desk", "Margin exception"),
          ("Credit / Approver", "Credit & budget")]
cw = 2.35; gap = 0.12; x0 = 0.6; y0 = 5.05; chh = 1.35
for i, (t, d) in enumerate(routes):
    px = x0 + i*(cw+gap)
    box(s, px, y0, cw, chh, fill=PANEL2)
    txt(s, px+0.22, y0+0.22, cw-0.4, 0.7, [[(t, 12, WHITE, True, False)]], line_spacing=0.95)
    txt(s, px+0.22, y0+0.9, cw-0.4, 0.4, [[(d, 10, YELLOW, False, True)]])

# ── Audit & governance ────────────────────────────────────────────────────────
s = slide()
header(s, "What we built", "Built-in audit & governance")
box(s, 0.6, 1.95, 6.0, 4.5, fill=PANEL)
txt(s, 0.9, 2.15, 5.5, 0.5, [[("Every decision is captured", 16, WHITE, True, False)]])
for i, t in enumerate(["What the agent decided automatically from master data",
                       "Why it paused for a CSR (the exception & rationale)",
                       "The exact action the CSR took (approve / pick / enter)",
                       "The final outcome applied to the order",
                       "Step-by-step audit trail per decision layer"]):
    yy = 2.75 + i*0.68
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(yy+0.05), Inches(0.16), Inches(0.16)); _set_fill(dot, YELLOW)
    txt(s, 1.25, yy, 5.1, 0.6, [[(t, 12, GREY, False, False)]], line_spacing=1.05)
box(s, 6.73, 1.95, 6.0, 4.5, fill=PANEL2)
box(s, 6.73, 1.95, 6.0, 0.12, fill=YELLOW, radius=False)
txt(s, 7.03, 2.2, 5.5, 0.4, [[("WHY IT MATTERS", 11, YELLOW, True, False)]])
for i, (t, d) in enumerate([("Governance", "Consistent, policy-aligned decisions across all CSRs."),
                            ("Traceability", "Full lineage of automated and human decisions."),
                            ("Continuous improvement", "Exception patterns reveal data & policy gaps to fix."),
                            ("Compliance-ready", "Defensible record for audits and disputes.")]):
    yy = 2.75 + i*0.9
    txt(s, 7.03, yy, 5.4, 0.4, [[(t, 13.5, WHITE, True, False)]])
    txt(s, 7.03, yy+0.36, 5.4, 0.5, [[(d, 11, GREY, False, False)]], line_spacing=1.0)

# ── Technology & deployment ───────────────────────────────────────────────────
s = slide()
header(s, "What we built", "Technology & deployment")
cols = [
    ("POC (today)", GREEN, ["Python orchestration engine", "Streamlit CSR workspace UI", "Master data as Excel workbooks",
                            "Mock integrations (email / ERP)", "Runs locally \u2014 fast to demo & iterate"]),
    ("Architecture principles", YELLOW, ["Modular agents & decision layers", "Resumable pipeline / state machine",
                            "Data-driven, not hard-coded rules", "Human-in-the-loop by design", "Audit-first for governance"]),
    ("Production path", BLUE, ["Connect to ERP / MDM / pricing systems", "Swap workbooks for live master data",
                            "LLM-assisted extraction & reasoning", "Enterprise auth, roles & SLAs", "Cloud deployment & monitoring"]),
]
cw = 3.95; gap = 0.13; x0 = 0.6; y0 = 1.95; chh = 4.5
for i, (t, col, items) in enumerate(cols):
    px = x0 + i*(cw+gap)
    box(s, px, y0, cw, chh, fill=PANEL)
    box(s, px, y0, cw, 0.12, fill=col, radius=False)
    txt(s, px+0.3, y0+0.35, cw-0.55, 0.5, [[(t, 15, WHITE, True, False)]])
    for j, it in enumerate(items):
        yy = y0 + 1.05 + j*0.68
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px+0.32), Inches(yy+0.05), Inches(0.14), Inches(0.14)); _set_fill(dot, col)
        txt(s, px+0.6, yy, cw-0.85, 0.6, [[(it, 11.5, GREY, False, False)]], line_spacing=1.0)

# ══════════════════════════════════════════════════════════════════════════════
# CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.333, 0.16, fill=YELLOW, radius=False)
box(s, 0, 7.34, 13.333, 0.16, fill=YELLOW, radius=False)
txt(s, 0.9, 2.4, 11.5, 0.5, [[("THANK YOU", 14, YELLOW, True, False)]])
txt(s, 0.9, 2.95, 11.5, 1.2, [[("Autonomous PO-to-Fulfillment Orchestration", 34, WHITE, True, False)]])
box(s, 0.94, 4.15, 3.0, 0.04, fill=YELLOW, radius=False)
txt(s, 0.9, 4.4, 11.4, 1.0,
    [[("A working accelerator today \u2014 a clear path to production.", 15, GREY, False, False)],
     [("Next: connect live systems, add LLM-assisted reasoning, and pilot with real production POs.", 13, DGREY, False, True)]],
    line_spacing=1.2)

# ── Save ──────────────────────────────────────────────────────────────────────
import os
out = os.path.join("demo", "PO-Fulfillment-Orchestration-Architecture.pptx")
prs.save(out)
print("Saved:", out, "slides:", len(prs.slides._sldIdLst))
